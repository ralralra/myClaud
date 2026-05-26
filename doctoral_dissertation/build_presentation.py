"""Build the 6-slide proceedings presentation as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

NAVY = RGBColor(0x1F, 0x36, 0x5C)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
TEXT = RGBColor(0x33, 0x33, 0x33)
SUBTEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_TXT = RGBColor(0xE8, 0xEE, 0xF7)

DIM_COLORS = {
    "문제 이해": RGBColor(0x2E, 0x86, 0xAB),
    "위임 인식": RGBColor(0x3E, 0x99, 0x53),
    "검증 판단": RGBColor(0xE6, 0x7E, 0x22),
    "책임 성찰": RGBColor(0x8E, 0x44, 0xAD),
}
DIM_BG = {k: RGBColor(*(min(255, c + 80) for c in (v[0], v[1], v[2]))) for k, v in {
    "문제 이해": (0x2E, 0x86, 0xAB),
    "위임 인식": (0x3E, 0x99, 0x53),
    "검증 판단": (0xE6, 0x7E, 0x22),
    "책임 성찰": (0x8E, 0x44, 0xAD),
}.items()}
DIM_BG_LIGHT = {
    "문제 이해": RGBColor(0xE3, 0xEF, 0xF6),
    "위임 인식": RGBColor(0xE3, 0xF1, 0xE6),
    "검증 판단": RGBColor(0xFB, 0xEC, 0xDC),
    "책임 성찰": RGBColor(0xEE, 0xE2, 0xF2),
}

FONT_KR = "Malgun Gothic"
FONT_EN = "Calibri"


def set_run(run, text, size=14, bold=False, color=TEXT, font=FONT_KR):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rpr = run._r.get_or_add_rPr()
    for tag in ("eastAsia", "cs"):
        existing = rpr.find(qn(f"a:{tag}"))
        if existing is not None:
            rpr.remove(existing)
    ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", font)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, font=font)
    return tb


def add_paragraph(tf, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
                  font=FONT_KR, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, font=font)
    return p


def add_filled_rect(slide, left, top, width, height, fill, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if not shadow:
        sppr = shape._element.spPr
        existing = sppr.find(qn("a:effectLst"))
        if existing is not None:
            sppr.remove(existing)
        etree.SubElement(sppr, qn("a:effectLst"))
    return shape


def add_rounded_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    sppr = shape._element.spPr
    existing = sppr.find(qn("a:effectLst"))
    if existing is not None:
        sppr.remove(existing)
    etree.SubElement(sppr, qn("a:effectLst"))
    return shape


def add_arrow(slide, left, top, width, height, fill=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_header(slide, prs, label, page_num, total=6):
    add_filled_rect(slide, 0, 0, prs.slide_width, Inches(0.45), NAVY)
    add_filled_rect(slide, 0, Inches(0.45), prs.slide_width, Inches(0.04), ACCENT)
    add_textbox(slide, Inches(0.4), Inches(0.06), Inches(10), Inches(0.35),
                label, size=14, bold=True, color=WHITE)
    add_textbox(slide, Inches(11.5), Inches(0.06), Inches(1.7), Inches(0.35),
                f"{page_num} / {total}", size=14, color=HEADER_TXT, align=PP_ALIGN.RIGHT)


def add_title(slide, text, top=Inches(0.65)):
    add_textbox(slide, Inches(0.5), top, Inches(12.3), Inches(0.7),
                text, size=30, bold=True, color=NAVY)


def add_footer(slide, prs, text="생성형 AI 활용 학습에서 사고 주도성과 메타인지적 조절 기제"):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
                text, size=10, color=SUBTEXT)


def style_cell(cell, text, size=12, bold=False, color=TEXT, bg=None,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_KR):
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = anchor
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, font=font)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---------------- Slide 1: 연구 배경·문제·연구문제 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, prs, "Ⅰ. 서론", 1)
    add_title(s, "연구 배경과 문제 의식")

    # Lead statement
    add_textbox(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5),
                "생성형 AI의 확산으로 학습 산출 방식이 근본적으로 변화하고 있다",
                size=18, color=TEXT)

    # Two contrast cards
    card_w, card_h = Inches(5.9), Inches(1.7)
    # Left: 변화
    add_rounded_rect(s, Inches(0.5), Inches(2.1), card_w, card_h, LIGHT_BG, line=NAVY)
    add_textbox(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(0.4),
                "학습 산출 방식의 변화", size=14, bold=True, color=NAVY)
    add_textbox(s, Inches(0.7), Inches(2.6), Inches(5.5), Inches(1.1),
                "글 · 보고서 · 발표자료 · 코드 등 다양한 학습 산출물을\n학습자가 직접 작성하지 않아도 AI가 즉시 생성한다.",
                size=13, color=TEXT)

    # Right: 위험
    add_rounded_rect(s, Inches(6.95), Inches(2.1), card_w, card_h,
                     RGBColor(0xFB, 0xEC, 0xEC), line=ACCENT)
    add_textbox(s, Inches(7.15), Inches(2.2), Inches(5.5), Inches(0.4),
                "사고 위임의 위험", size=14, bold=True, color=ACCENT)
    add_textbox(s, Inches(7.15), Inches(2.6), Inches(5.5), Inches(1.1),
                "산출물은 완성되어도 학습자의 사고는 빠질 수 있다.\nAI 고스트라이터 효과(Draxler, 2024), 인지적 오프로딩을\n매개로 한 비판적 사고 약화(Gerlich, 2025) 보고.",
                size=12, color=TEXT)

    # Research questions header
    add_textbox(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.4),
                "연구문제", size=16, bold=True, color=NAVY)

    rqs = [
        ("RQ 1", "AI 활용 학습에서 학습자의 사고 주도성은 왜 약화될 수 있는가?"),
        ("RQ 2", "사고 주도성 유지를 위해 메타인지적 조절 기제가 필요한 이유는 무엇인가?"),
        ("RQ 3", "메타인지적 조절 기제는 어떤 방향으로 설계될 필요가 있는가?"),
    ]
    rq_w = Inches(4.05)
    for i, (tag, body) in enumerate(rqs):
        left = Inches(0.5 + i * 4.2)
        add_rounded_rect(s, left, Inches(4.55), rq_w, Inches(2.15), WHITE, line=NAVY)
        add_filled_rect(s, left, Inches(4.55), rq_w, Inches(0.55), NAVY)
        add_textbox(s, left + Inches(0.2), Inches(4.6), rq_w - Inches(0.2),
                    Inches(0.45), tag, size=15, bold=True, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, left + Inches(0.2), Inches(5.25), rq_w - Inches(0.4),
                    Inches(1.4), body, size=13, color=TEXT)

    add_footer(s, prs)

    # ---------------- Slide 2: 이론적 배경 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, prs, "Ⅱ. 이론적 배경", 2)
    add_title(s, "이론적 배경 — 인지적 오프로딩 · 메타인지 · 자기조절학습")

    # Three concept boxes with arrows
    box_w = Inches(3.7)
    box_h = Inches(2.7)
    y_box = Inches(1.65)
    gap = Inches(0.4)
    x_starts = [Inches(0.5), Inches(0.5 + 3.7 + 0.4), Inches(0.5 + (3.7 + 0.4) * 2)]

    items = [
        ("인지적 오프로딩",
         "정보처리 요구를 외부 도구·환경에\n분산하는 행위. 효율적 전략일 수 있으나\n메타인지적 평가의 오류로 인해\n학습 사고까지 위임될 위험이 있다.",
         "Risko & Gilbert (2016)"),
        ("메타인지",
         "인지에 대한 지식과 조절.\n계획·점검·평가의 조절 과정이\n학습 사고를 점검·수정하는\n핵심 기제로 작동한다.",
         "Flavell (1979); Schraw & Dennison (1994)"),
        ("자기조절학습",
         "사전숙고 → 수행 통제 → 자기성찰의\n순환적 학습 모형.\n메타인지 조절은 학습 전·중·후\n전 과정에 분산 배치되어야 한다.",
         "Zimmerman (2002)"),
    ]
    for i, (title, body, cite) in enumerate(items):
        add_rounded_rect(s, x_starts[i], y_box, box_w, box_h, WHITE, line=NAVY)
        add_filled_rect(s, x_starts[i], y_box, box_w, Inches(0.55), NAVY)
        add_textbox(s, x_starts[i] + Inches(0.2), y_box + Inches(0.07),
                    box_w - Inches(0.2), Inches(0.45),
                    title, size=15, bold=True, color=WHITE)
        add_textbox(s, x_starts[i] + Inches(0.2), y_box + Inches(0.7),
                    box_w - Inches(0.4), Inches(1.6),
                    body, size=12, color=TEXT)
        add_textbox(s, x_starts[i] + Inches(0.2), y_box + Inches(2.3),
                    box_w - Inches(0.4), Inches(0.35),
                    cite, size=10, color=SUBTEXT)
        if i < 2:
            arrow_x = x_starts[i] + box_w + Inches(0.02)
            add_arrow(s, arrow_x, y_box + Inches(1.2), Inches(0.36), Inches(0.3), NAVY)

    # Bottom takeaway
    add_rounded_rect(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.95),
                     LIGHT_BG, line=NAVY)
    add_textbox(s, Inches(0.8), Inches(4.95), Inches(11.8), Inches(0.45),
                "본 연구의 이론적 함의", size=15, bold=True, color=NAVY)
    add_textbox(s, Inches(0.8), Inches(5.4), Inches(11.8), Inches(1.4),
                "AI 활용 학습에서는 메타인지적 평가의 오류 가능성이 더 크게 작동한다.\n"
                "산출물이 매끄럽게 제시되어 이해 수준을 과대평가하기 쉽고,\n"
                "AI 산출물의 정확성을 과신하기 쉬우며, 위임의 학습 손실이 즉시 드러나지 않는다.\n"
                "따라서 메타인지 조절은 사후 성찰이 아닌 학습 전 과정의 설계 요소가 되어야 한다.",
                size=13, color=TEXT)

    add_footer(s, prs)

    # ---------------- Slide 3: 사고 주도성 개념 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, prs, "Ⅲ. 사고 주도성", 3)
    add_title(s, "핵심 개념 — 사고 주도성 (Thinking Agency)")

    # Definition card
    add_rounded_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.85),
                     LIGHT_BG, line=NAVY)
    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.8), Inches(0.45),
                "정의", size=14, bold=True, color=NAVY)
    add_textbox(s, Inches(0.8), Inches(2.05), Inches(11.8), Inches(1.3),
                "생성형 AI를 활용하는 학습 과정에서 학습자가 문제 이해, 질문 생성,\n"
                "산출물 검토, 최종 판단, 책임 인식을 AI에 전적으로 위임하지 않고\n"
                "유지하는 상태. AI 사용량의 문제가 아니라 AI 활용 과정에서의\n"
                "‘학습자 위치’의 문제이다.",
                size=14, color=TEXT)

    # Three distinctions
    add_textbox(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.4),
                "관련 개념과의 구분", size=15, bold=True, color=NAVY)

    distinctions = [
        ("자기조절학습",
         "학습 목표 달성을 위한 동기·전략·행동의 포괄적 조절.\n"
         "사고 주도성은 그중 ‘AI–학습자 사이 사고 분담’에 초점."),
        ("인식론적 행위주체성",
         "학습자가 지식 생산의 주체로 남아 있는가에 대한 거시적 질문.\n"
         "사고 주도성은 AI 상호작용 단위의 미시적 구현."),
        ("AI 리터러시",
         "AI를 비판·활용·소통하는 역량.\n"
         "사고 주도성은 그 역량이 학습자 위치로 실제 유지되는가를 다룸."),
    ]
    card_w = Inches(4.05)
    for i, (title, body) in enumerate(distinctions):
        left = Inches(0.5 + i * 4.2)
        add_rounded_rect(s, left, Inches(4.15), card_w, Inches(2.55), WHITE, line=NAVY)
        add_filled_rect(s, left, Inches(4.15), card_w, Inches(0.55), NAVY)
        add_textbox(s, left + Inches(0.2), Inches(4.21), card_w - Inches(0.2),
                    Inches(0.45), title, size=14, bold=True, color=WHITE)
        add_textbox(s, left + Inches(0.2), Inches(4.85), card_w - Inches(0.4),
                    Inches(1.8), body, size=12, color=TEXT)

    add_footer(s, prs)

    # ---------------- Slide 4: 사고 주도성의 네 차원 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, prs, "Ⅲ. 사고 주도성", 4)
    add_title(s, "사고 주도성의 네 차원")

    rows = [
        ("문제 이해", "AI에게 묻기 전 학습자가 과제의 목적과\n문제를 이해하는 정도", "‘무엇을 묻는가’의 주체"),
        ("위임 인식", "AI가 수행한 부분과 자신이 수행한\n부분을 구분하는 정도", "사고 분담에 대한 자각"),
        ("검증 판단", "AI 산출물의 정확성, 근거, 맥락\n적합성을 점검하는 정도", "산출물에 대한 판단의 주체"),
        ("책임 성찰", "최종 결과물에 대한 자신의 이해와\n책임을 설명하는 정도", "결과물에 대한 책임의 주체"),
    ]
    col_widths = [Inches(2.2), Inches(6.0), Inches(4.1)]
    n_cols = 3
    n_rows = len(rows) + 1
    table_left = Inches(0.5)
    table_top = Inches(1.7)
    row_h = Inches(1.05)
    header_h = Inches(0.55)
    total_w = sum(col_widths, Inches(0))

    table_shape = s.shapes.add_table(n_rows, n_cols, table_left, table_top, total_w,
                                     header_h + row_h * (n_rows - 1))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    table.rows[0].height = header_h
    for i in range(1, n_rows):
        table.rows[i].height = row_h

    headers = ["차원", "의미", "학습자 사고 위치의 핵심"]
    for i, h in enumerate(headers):
        style_cell(table.cell(0, i), h, size=14, bold=True, color=WHITE,
                   bg=NAVY, align=PP_ALIGN.CENTER)

    for r, (dim, meaning, focus) in enumerate(rows, start=1):
        dim_color = DIM_COLORS[dim]
        dim_bg = DIM_BG_LIGHT[dim]
        style_cell(table.cell(r, 0), dim, size=14, bold=True, color=WHITE,
                   bg=dim_color, align=PP_ALIGN.CENTER)
        style_cell(table.cell(r, 1), meaning, size=12, bg=WHITE)
        style_cell(table.cell(r, 2), focus, size=12, bg=dim_bg, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.45),
                "→ 네 차원은 독립된 단계가 아니라, AI 활용 학습 과정에서 시점에 따라 강조점이 달라지는 구성 요소이다.",
                size=12, color=SUBTEXT)

    add_footer(s, prs)

    # ---------------- Slide 5: 메타인지적 조절 기제 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, prs, "Ⅳ. 조절 기제와 설계원리", 5)
    add_title(s, "메타인지적 조절 기제 — 활용 전 · 중 · 후")

    # Top timeline
    timeline_top = Inches(1.6)
    timeline_h = Inches(0.4)
    add_textbox(s, Inches(0.5), timeline_top, Inches(1.0), timeline_h,
                "시점", size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    stages = [("활용 전", Inches(1.6), Inches(2.7)),
              ("활용 중", Inches(4.5), Inches(2.7)),
              ("활용 후", Inches(7.4), Inches(2.7)),
              ("최종", Inches(10.3), Inches(2.5))]
    for label, left, width in stages:
        add_rounded_rect(s, left, timeline_top, width, timeline_h, NAVY)
        add_textbox(s, left, timeline_top + Inches(0.04), width, timeline_h - Inches(0.04),
                    label, size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("문제 이해", "활용 전 자기점검",
         "나는 이 과제에서\n무엇을 알고, 무엇을 모르는가?",
         "사전 문제·가설 작성, 질문 작성"),
        ("위임 인식", "활용 중 위임 점검",
         "내가 직접 생각한 부분과\nAI가 대신한 부분은 무엇인가?",
         "AI 위임표, 프롬프트 로그"),
        ("검증 판단", "활용 후 검증 점검",
         "이 결과를 믿을 수 있는\n근거는 무엇인가?",
         "오류검증표, 출처 확인, 코드·문장 검토"),
        ("책임 성찰", "최종 책임 성찰",
         "나는 이 결과물을 내 언어로\n설명하고 책임질 수 있는가?",
         "성찰문, 발표 전 자기점검,\nAI 활용 범위 진술"),
    ]
    col_widths = [Inches(1.6), Inches(2.4), Inches(4.3), Inches(4.0)]
    n_cols = 4
    n_rows = len(rows) + 1
    table_left = Inches(0.5)
    table_top = Inches(2.3)
    header_h = Inches(0.55)
    row_h = Inches(0.95)
    total_w = sum(col_widths, Inches(0))

    table_shape = s.shapes.add_table(n_rows, n_cols, table_left, table_top, total_w,
                                     header_h + row_h * (n_rows - 1))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    table.rows[0].height = header_h
    for i in range(1, n_rows):
        table.rows[i].height = row_h

    headers = ["사고 주도성 차원", "조절 기제", "핵심 점검 질문", "수업 적용 및 확인 자료"]
    for i, h in enumerate(headers):
        style_cell(table.cell(0, i), h, size=13, bold=True, color=WHITE,
                   bg=NAVY, align=PP_ALIGN.CENTER)

    for r, (dim, mech, q, mat) in enumerate(rows, start=1):
        dim_color = DIM_COLORS[dim]
        dim_bg = DIM_BG_LIGHT[dim]
        style_cell(table.cell(r, 0), dim, size=13, bold=True, color=WHITE,
                   bg=dim_color, align=PP_ALIGN.CENTER)
        style_cell(table.cell(r, 1), mech, size=12, bold=True, color=dim_color, bg=dim_bg)
        style_cell(table.cell(r, 2), q, size=11, bg=WHITE)
        style_cell(table.cell(r, 3), mat, size=11, bg=WHITE)

    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.45),
                "→ Schraw·Dennison(1994)의 계획·점검·평가와 Zimmerman(2002)의 사전숙고–수행–자기성찰 순환에 대응.",
                size=12, color=SUBTEXT)

    add_footer(s, prs)

    # ---------------- Slide 6: 결론 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, prs, "Ⅴ. 결론", 6)
    add_title(s, "결론과 후속 연구")

    # Three column summary
    col_w = Inches(4.05)
    col_h = Inches(2.6)
    col_top = Inches(1.55)
    titles = [
        ("결론", NAVY,
         "AI 활용 학습의 설계 논의를\n‘산출물 중심’에서\n‘사고 과정 중심’으로 전환.\n\n"
         "사고 주도성 4차원과\n메타인지적 조절 기제 4종을 제안."),
        ("후속 연구 절차", ACCENT,
         "① 조절 기제를 문항·활동지로 개발\n\n"
         "② 전문가 검토를 통한\n   타당성 확인 (델파이 등)\n\n"
         "③ 실제 학습 현장 적용·효과 검증"),
        ("한계", SUBTEXT,
         "본 연구는 문헌 고찰에 기반한\n개념적 논의로,\n"
         "사고 주도성과 조절 기제는\n아직 경험적으로 검증되지 않았다.\n\n"
         "후속 실증 연구가 요구된다."),
    ]
    for i, (title, color, body) in enumerate(titles):
        left = Inches(0.5 + i * 4.2)
        add_rounded_rect(s, left, col_top, col_w, col_h, WHITE, line=color)
        add_filled_rect(s, left, col_top, col_w, Inches(0.55), color)
        add_textbox(s, left + Inches(0.2), col_top + Inches(0.07),
                    col_w - Inches(0.2), Inches(0.45),
                    title, size=15, bold=True, color=WHITE)
        add_textbox(s, left + Inches(0.2), col_top + Inches(0.7),
                    col_w - Inches(0.4), col_h - Inches(0.8),
                    body, size=12, color=TEXT)

    # Closing one-liner
    add_filled_rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.0), NAVY)
    add_filled_rect(s, Inches(0.5), Inches(4.6), Inches(0.18), Inches(2.0), ACCENT)
    add_textbox(s, Inches(0.9), Inches(4.75), Inches(11.6), Inches(0.45),
                "한 줄 메시지", size=13, bold=True, color=HEADER_TXT)
    add_textbox(s, Inches(0.9), Inches(5.2), Inches(11.6), Inches(1.3),
                "AI 시대의 교육은 AI를 능숙하게 사용하는 학습자를 기르는 데서 나아가,\n"
                "학습자가 ‘사고의 주체’로 남도록 설계되어야 한다.",
                size=20, bold=True, color=WHITE)

    add_footer(s, prs)

    out_path = "doctoral_dissertation/presentation.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    build()
